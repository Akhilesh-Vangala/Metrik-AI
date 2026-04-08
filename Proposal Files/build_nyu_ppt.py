"""
build_nyu_ppt.py — Metrik AI · DS-GA 1019 Proposal · NYU Bold Official Template
7 slides: Title · Motivation+Problem · Dataset · ML Pipeline · Optimization
          Literature · Goals
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

TEMPLATE = "Copy of NYU Presentation (Bold) Official Template.pptx"
OUTPUT   = "MetrikAI_NYU_Proposal.pptx"

# ── NYU Brand ─────────────────────────────────────────────────────────────────
NYU_PURPLE  = "57068C"   # dark purple  – backgrounds, bold titles
NYU_VIOLET  = "8900E1"   # bright violet – accents, dividers, highlights
WHITE       = "FFFFFF"
DARK        = "1A1A2E"   # very dark navy for text
MID         = "555555"   # medium gray body text
CARD_BG     = "F5F1FB"   # light lavender card

FONT_H = "Montserrat Black"    # headings
FONT_B = "Montserrat"          # body

# ── Canvas constants ──────────────────────────────────────────────────────────
W  = Inches(10)
H  = Inches(5.625)

# Safe content zone — NYU logo sits at (0.445", 4.964") on white slides
# and the top bar occupies y=0 to 0.336"
Y_TITLE   = Inches(0.487)    # title placeholder top
Y_CONTENT = Inches(1.18)     # first content row after title + divider
Y_MAX     = Inches(4.82)     # bottom content boundary (above NYU logo)
X_L       = Inches(0.34)     # left margin
X_R       = Inches(9.65)     # right boundary
SAFE_W    = Inches(9.30)     # usable width


# ══════════════════════════════════════════════════════════════════════════════
# XML helpers
# ══════════════════════════════════════════════════════════════════════════════
ANS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _p(lines, ns=ANS):
    """Build lxml <a:p> elements for a txBody from line-dicts."""
    elems = []
    for ln in lines:
        p = etree.Element(f"{{{ns}}}p")
        pPr = etree.SubElement(p, f"{{{ns}}}pPr")
        pPr.set("algn",   ln.get("align", "l"))
        pPr.set("indent", "0")
        pPr.set("marL",   "0")
        sb  = etree.SubElement(pPr, f"{{{ns}}}spcBef")
        spc = etree.SubElement(sb,  f"{{{ns}}}spcPts")
        spc.set("val", str(ln.get("sp", 0)))
        etree.SubElement(pPr, f"{{{ns}}}buNone")

        txt = ln.get("text", "")
        if txt:
            r   = etree.SubElement(p, f"{{{ns}}}r")
            rPr = etree.SubElement(r, f"{{{ns}}}rPr")
            rPr.set("lang", "en-US")
            rPr.set("sz",   str(int(ln.get("sz", 14) * 100)))
            rPr.set("b",    "1" if ln.get("b", False)  else "0")
            rPr.set("i",    "1" if ln.get("i", False)  else "0")
            rPr.set("dirty","0")
            fl  = etree.SubElement(rPr, f"{{{ns}}}solidFill")
            clr = etree.SubElement(fl,  f"{{{ns}}}srgbClr")
            clr.set("val", ln.get("c", MID))
            lat = etree.SubElement(rPr, f"{{{ns}}}latin")
            lat.set("typeface", ln.get("f", FONT_B))
            t = etree.SubElement(r, f"{{{ns}}}t")
            t.text = txt
        elems.append(p)
    return elems


def set_ph(ph, lines):
    """Replace all paragraphs in a placeholder with given lines."""
    tf    = ph.text_frame
    txBody= tf._txBody
    ns    = ANS
    for old in txBody.findall(f"{{{ns}}}p"):
        txBody.remove(old)
    bodyPr = txBody.find(f"{{{ns}}}bodyPr")
    if bodyPr is not None:
        bodyPr.set("wrap", "square")
    for pel in _p(lines):
        txBody.append(pel)


def tb(slide, lines, left, top, width, height):
    """Add a textbox and fill with lines."""
    box   = slide.shapes.add_textbox(left, top, width, height)
    tf    = box.text_frame
    tf.word_wrap = True
    txBody= tf._txBody
    ns    = ANS
    for old in txBody.findall(f"{{{ns}}}p"):
        txBody.remove(old)
    bodyPr = txBody.find(f"{{{ns}}}bodyPr")
    if bodyPr is not None:
        bodyPr.set("wrap", "square")
        bodyPr.set("anchor", "t")
    for pel in _p(lines):
        txBody.append(pel)
    return box


def rect(slide, left, top, w, h, fill, border=None):
    sh = slide.shapes.add_shape(1, left, top, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor.from_string(fill)
    if border:
        sh.line.color.rgb = RGBColor.from_string(border)
        sh.line.width     = Pt(0.75)
    else:
        sh.line.fill.background()
    return sh


# ── Shorthand line constructors ───────────────────────────────────────────────
def L(text, sz=13, b=False, c=MID, f=FONT_B, i=False, align="l", sp=0):
    return dict(text=text, sz=sz, b=b, c=c, f=f, i=i, align=align, sp=sp)

def H1(text, c=NYU_PURPLE, sz=22): return L(text, sz=sz, b=True, c=c, f=FONT_H)
def H2(text, c=NYU_PURPLE, sz=14): return L(text, sz=sz, b=True, c=c, f=FONT_H)
def H3(text, c=NYU_VIOLET, sz=12): return L(text, sz=sz, b=True, c=c, f=FONT_H)
def BD(text, sz=11.5, c=MID):      return L(text, sz=sz, c=c)
def BU(text, sz=11.5, c=MID):      return L(f"\u25b8  {text}", sz=sz, c=c)
def SP(pts=250):                    return L("", sz=6,   sp=pts)
def WH(text, sz=13, al="ctr"):     return L(text, sz=sz, b=True, c=WHITE, f=FONT_H, align=al)


# ── Remove all slides keeping master/layouts ──────────────────────────────────
def clear_slides(prs):
    lst = prs.slides._sldIdLst
    rIds = [
        sid.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        for sid in list(lst)
    ]
    for sid in list(lst):
        lst.remove(sid)
    for rId in rIds:
        try: prs.part.drop_rel(rId)
        except: pass


# ── Common chrome ─────────────────────────────────────────────────────────────
def divider(slide, y=Inches(1.10), x=X_L, w=SAFE_W):
    rect(slide, x, y, w, Inches(0.04), NYU_VIOLET)


def section_title(slide, text, sub=None):
    """Set title placeholder and optionally add a subtitle."""
    ph = slide.shapes.placeholders[0]
    if sub:
        set_ph(ph, [H1(text, sz=18), L(sub, sz=12, c=NYU_VIOLET, f=FONT_B, b=True)])
    else:
        set_ph(ph, [H1(text, sz=18)])
    divider(slide)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

# ── 1. TITLE ──────────────────────────────────────────────────────────────────
def s1_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])   # TITLE – purple bg

    set_ph(slide.shapes.placeholders[0], [
        L("METRIK AI", sz=54, b=True, c=WHITE, f=FONT_H),
    ])
    set_ph(slide.shapes.placeholders[1], [
        L("Building Energy Consumption Prediction & Optimization",
          sz=20, c=WHITE, f=FONT_B),
    ])
    set_ph(slide.shapes.placeholders[2], [
        L("DS-GA 1019  ·  Advanced Python for Data Science  ·  Spring 2026",
          sz=13, c=WHITE, f=FONT_B),
        SP(200),
        L("NYU Center for Data Science  ·  March 4, 2026",
          sz=11, c=WHITE, f=FONT_B, i=True),
    ])

    # Bottom violet strip
    rect(slide, Inches(0), H - Inches(0.10), W, Inches(0.10), NYU_VIOLET)

    # Three stat tiles — right side, well above strip and placeholders
    for i, (val, lbl) in enumerate([("53.6M","meter readings"),
                                     ("1,636","buildings"),
                                     ("19",   "sites")]):
        bx = Inches(6.1) + Inches(i * 1.28)
        by = Inches(3.88)
        rect(slide, bx, by, Inches(1.14), Inches(0.90), NYU_VIOLET)
        tb(slide, [WH(val, sz=21), SP(80), WH(lbl, sz=9)],
           bx, by + Inches(0.06), Inches(1.14), Inches(0.82))


# ── 2. REAL-WORLD MOTIVATION + PROBLEM STATEMENT ─────────────────────────────
def s2_motivation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[7])

    section_title(slide, "REAL-WORLD MOTIVATION  &  PROBLEM STATEMENT")

    # Exact column math: COL1_W + GAP + COL2_W = SAFE_W
    COL1_W = Inches(4.55)
    GAP    = Inches(0.20)
    COL2_W = SAFE_W - COL1_W - GAP        # = 4.55"
    COL2_X = X_L + COL1_W + GAP
    CONTENT_H = Y_MAX - Y_CONTENT          # = 3.64"

    # ── Column 1: Motivation ──────────────────────────────────────────────────
    HDR_H  = Inches(0.35)
    rect(slide, X_L, Y_CONTENT, COL1_W, HDR_H, NYU_PURPLE)
    tb(slide, [WH("REAL-WORLD MOTIVATION", sz=11)],
       X_L, Y_CONTENT + Inches(0.03), COL1_W, HDR_H)

    # Stat tiles: 2 × 2, each tile fills exactly half col width
    TILE_W  = (COL1_W - Inches(0.08)) / 2   # small gap between tiles
    TILE_H  = Inches(0.90)
    TILE_GAP= Inches(0.08)
    stat_y  = Y_CONTENT + HDR_H + Inches(0.06)
    for i, (val, lbl, col) in enumerate([
        ("40%",   "of U.S. energy consumed by buildings",    NYU_PURPLE),
        ("28%",   "of global CO₂ from building operations",  NYU_VIOLET),
        ("20–30%","of building energy typically wasted",     NYU_PURPLE),
        ("$130B+","annual cost of building energy waste",    NYU_VIOLET),
    ]):
        c_i, r_i = i % 2, i // 2
        bx = X_L + c_i * (TILE_W + TILE_GAP)
        by = stat_y + r_i * (TILE_H + Inches(0.06))
        rect(slide, bx, by, TILE_W, TILE_H, col)
        tb(slide, [WH(val, sz=20), SP(50), WH(lbl, sz=9)],
           bx, by + Inches(0.05), TILE_W, TILE_H - Inches(0.05))

    # "Who is affected" — fills remaining space exactly to Y_MAX
    who_y   = stat_y + 2 * (TILE_H + Inches(0.06)) + Inches(0.08)
    who_h   = Y_MAX - who_y
    rect(slide, X_L, who_y, COL1_W, Inches(0.03), NYU_VIOLET)
    tb(slide, [
        H3("Who is affected?", sz=11),
        SP(150),
        BU("Facility managers & building operators", sz=11),
        BU("Utilities & grid demand planners",       sz=11),
        BU("ESG / sustainability teams",             sz=11),
        BU("Energy services companies (ESCOs)",      sz=11),
    ], X_L, who_y + Inches(0.08), COL1_W, who_h - Inches(0.08))

    # Vertical divider
    rect(slide, COL2_X - GAP / 2 - Inches(0.02), Y_CONTENT,
         Inches(0.04), CONTENT_H, NYU_VIOLET)

    # ── Column 2: Problem Statement ──────────────────────────────────────────
    rect(slide, COL2_X, Y_CONTENT, COL2_W, HDR_H, NYU_VIOLET)
    tb(slide, [WH("PROBLEM STATEMENT", sz=11)],
       COL2_X, Y_CONTENT + Inches(0.03), COL2_W, HDR_H)

    # Formal statement callout
    CALLOUT_H = Inches(0.72)
    ps_y = Y_CONTENT + HDR_H + Inches(0.06)
    rect(slide, COL2_X, ps_y, COL2_W, CALLOUT_H, CARD_BG, NYU_VIOLET)
    tb(slide, [
        L("Given hourly meter readings from 1,636 non-residential buildings", sz=11, i=True, c=DARK),
        L("across 19 sites over two full years — predict next-hour energy",   sz=11, i=True, c=DARK),
        L("consumption per meter at scale, without exceeding RAM limits.",     sz=11, i=True, c=DARK, b=True),
    ], COL2_X + Inches(0.1), ps_y + Inches(0.06), COL2_W - Inches(0.15), CALLOUT_H - Inches(0.08))

    # Three components — divide remaining height exactly
    remaining = Y_MAX - (ps_y + CALLOUT_H + Inches(0.06))
    COMP_H    = remaining / 3
    comps = [
        ("FORECASTING",      NYU_PURPLE,
         "Predict next-hour meter-level consumption\n→ Enables demand response & pre-conditioning"),
        ("ANOMALY DETECTION",NYU_VIOLET,
         "Flag residuals exceeding statistical threshold\n→ Detects faults, outages & billing errors"),
        ("DECISION SUPPORT", NYU_PURPLE,
         "Rank buildings & meters by anomaly severity\n→ Tells operators exactly where to act first"),
    ]
    COMP_HDR = Inches(0.30)
    for j, (title, col, body) in enumerate(comps):
        cy = ps_y + CALLOUT_H + Inches(0.06) + j * COMP_H
        rect(slide, COL2_X, cy, COL2_W, COMP_HDR, col)
        tb(slide, [WH(title, sz=10)],
           COL2_X, cy + Inches(0.02), COL2_W, COMP_HDR)
        rect(slide, COL2_X, cy + COMP_HDR, COL2_W, COMP_H - COMP_HDR - Inches(0.03), CARD_BG)
        tb(slide, [BD(body, sz=11)],
           COL2_X + Inches(0.1), cy + COMP_HDR + Inches(0.05),
           COL2_W - Inches(0.15), COMP_H - COMP_HDR - Inches(0.08))


# ── 3. DATASET ────────────────────────────────────────────────────────────────
def s3_dataset(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[7])

    section_title(slide, "DATASET  —  ASHRAE GREAT ENERGY PREDICTOR III")

    COL1_W = Inches(4.55)
    GAP    = Inches(0.20)
    COL2_W = SAFE_W - COL1_W - GAP        # = 4.55"
    COL2_X = X_L + COL1_W + GAP
    CONTENT_H = Y_MAX - Y_CONTENT

    rect(slide, COL2_X - Inches(0.12), Y_CONTENT, Inches(0.04), CONTENT_H, NYU_VIOLET)

    # ── LEFT: 6 stat tiles (3×2) + files ─────────────────────────────────────
    TILE_W  = (COL1_W - Inches(0.10)) / 3    # exact third minus gaps
    TILE_H  = Inches(0.68)
    TILE_GX = Inches(0.05)
    TILE_GY = Inches(0.06)
    for i, (val, lbl, col) in enumerate([
        ("53.6M","Hourly training rows",       NYU_PURPLE),
        ("1,636","Non-residential buildings",  NYU_VIOLET),
        ("3,053","Energy meters",              NYU_PURPLE),
        ("19",   "Sites — NA & EU",            NYU_VIOLET),
        ("2 yrs","Jan 2016 – Dec 2017",        NYU_PURPLE),
        ("4",    "Meter types",                NYU_VIOLET),
    ]):
        ci, ri = i % 3, i // 3
        bx = X_L + ci * (TILE_W + TILE_GX)
        by = Y_CONTENT + ri * (TILE_H + TILE_GY)
        rect(slide, bx, by, TILE_W, TILE_H, col)
        tb(slide, [WH(val, sz=16), SP(30), WH(lbl, sz=8.5)],
           bx, by + Inches(0.04), TILE_W, TILE_H - Inches(0.04))

    # Files table fills remaining height to Y_MAX exactly
    tiles_bottom = Y_CONTENT + 2 * (TILE_H + TILE_GY)
    files_hdr_y  = tiles_bottom + Inches(0.08)
    tb(slide, [H3("DATASET FILES", sz=11)], X_L, files_hdr_y, COL1_W, Inches(0.26))

    files = [
        ("train.csv",             "53.6M rows · meter, timestamp, reading (target)"),
        ("building_metadata.csv", "1,636 rows · site, use type, sq ft, year built"),
        ("weather_train.csv",     "Hourly weather per site · temp, dew, cloud, wind"),
        ("test.csv",              "Competition test set · same schema, no target"),
    ]
    FILE_START = files_hdr_y + Inches(0.28)
    FILE_H     = (Y_MAX - FILE_START) / len(files)
    for i, (fn, desc) in enumerate(files):
        fy = FILE_START + i * FILE_H
        bg = CARD_BG if i % 2 == 0 else WHITE
        rect(slide, X_L, fy, COL1_W, FILE_H - Inches(0.02), bg, NYU_VIOLET)
        tb(slide, [L(fn,   sz=10.5, b=True,  c=NYU_VIOLET)],
           X_L + Inches(0.08), fy + Inches(0.04), Inches(1.72), FILE_H - Inches(0.06))
        tb(slide, [L(desc, sz=10.5, b=False, c=MID)],
           X_L + Inches(1.82), fy + Inches(0.05), COL1_W - Inches(1.88), FILE_H - Inches(0.06))

    # ── RIGHT: Target + Features + Significance ───────────────────────────────
    # Target callout
    tb(slide, [H3("TARGET VARIABLE", sz=11)], COL2_X, Y_CONTENT, COL2_W, Inches(0.26))
    rect(slide, COL2_X, Y_CONTENT + Inches(0.27), COL2_W, Inches(0.38), CARD_BG, NYU_VIOLET)
    tb(slide, [L("meter_reading  (float, kWh / kBTU)  —  continuous regression target",
               sz=11, i=True, c=DARK)],
       COL2_X + Inches(0.1), Y_CONTENT + Inches(0.30), COL2_W - Inches(0.15), Inches(0.32))

    # Features + Significance — fill to Y_MAX
    tb(slide, [
        H3("ENGINEERED FEATURES", sz=11),
        SP(150),
        BU("Temporal:  hour-of-day, day-of-week, month, holiday, weekend",     sz=11),
        BU("Lag:  same-hour-yesterday (lag=24), same-day-last-week (lag=168)",  sz=11),
        BU("Rolling:  24h & 168h mean — strictly past-only, leakage-safe",     sz=11),
        BU("Building:  log floor area, use type (16 categories), building age", sz=11),
        BU("Weather:  air temperature, dew point, rolling_temp_24h",            sz=11),
        SP(350),
        H3("SIGNIFICANCE", sz=11),
        SP(150),
        BU("Largest open non-residential energy benchmark (Miller et al. 2020)", sz=11),
        BU("Real operational data — 19 sites across North America & Europe",     sz=11),
        BU("Referenced in 100+ papers · IEEE Power & Energy Society benchmark",  sz=11),
        BU("53.6M rows demand every performance optimization in this project",   sz=11, c=NYU_VIOLET),
    ], COL2_X, Y_CONTENT + Inches(0.72), COL2_W, Y_MAX - Y_CONTENT - Inches(0.72))


# ── 4. ML PIPELINE (FLOWCHART) ────────────────────────────────────────────────
def s4_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[7])

    section_title(slide, "END-TO-END MACHINE LEARNING PIPELINE")

    # ── Flowchart row 1: 6 process boxes ─────────────────────────────────────
    steps = [
        ("RAW DATA",         NYU_PURPLE, "ASHRAE\n53.6M rows\n2 yrs · 19 sites"),
        ("OUT-OF-CORE\nLOAD",NYU_VIOLET, "Chunked reads\n1–2M rows/chunk\n<4 GB peak RAM"),
        ("FEATURE\nENGINEER",NYU_PURPLE, "Lag · Rolling\nTemporal · Weather\nBuilding metadata"),
        ("LIGHTGBM\nMODEL",  NYU_VIOLET, "Baseline vs.\nGBDT model\nRMSE / CV-RMSE"),
        ("ANOMALY\nSCORING", NYU_PURPLE, "Residuals\nModified Z-score\nPer-meter flags"),
        ("DECISION\nSUPPORT",NYU_VIOLET, "Ranked audit\nlist by severity\nCSV / JSON out"),
    ]

    N      = len(steps)
    BOX_W  = Inches(1.38)
    BOX_H  = Inches(0.60)
    ARR_W  = Inches(0.17)
    TOTAL  = N * BOX_W + (N - 1) * ARR_W
    START_X= (W - TOTAL) / 2
    BOX_Y  = Inches(1.22)
    DESC_Y = BOX_Y + BOX_H + Inches(0.10)
    DESC_H = Inches(0.78)

    for i, (label, col, desc) in enumerate(steps):
        bx = START_X + Inches(i) * (BOX_W + ARR_W)

        # Box
        rect(slide, bx, BOX_Y, BOX_W, BOX_H, col)
        tb(slide, [WH(label, sz=10)],
           bx, BOX_Y + Inches(0.04), BOX_W, BOX_H - Inches(0.04))

        # Arrow (→) between boxes
        if i < N - 1:
            ax = bx + BOX_W
            tb(slide, [L("→", sz=16, b=True, c=NYU_VIOLET, align="ctr")],
               ax, BOX_Y + Inches(0.14), ARR_W, Inches(0.32))

        # Step number circle
        rect(slide, bx + BOX_W / 2 - Inches(0.13), BOX_Y - Inches(0.24),
             Inches(0.26), Inches(0.24), col)
        tb(slide, [WH(str(i + 1), sz=10, al="ctr")],
           bx + BOX_W / 2 - Inches(0.13), BOX_Y - Inches(0.24),
           Inches(0.26), Inches(0.24))

        # Description card
        rect(slide, bx, DESC_Y, BOX_W, DESC_H, CARD_BG, col)
        tb(slide, [L(desc, sz=10, c=MID, align="ctr")],
           bx, DESC_Y + Inches(0.06), BOX_W, DESC_H - Inches(0.08))

    # ── Three technical pillars below ─────────────────────────────────────────
    pil_y   = DESC_Y + DESC_H + Inches(0.16)
    pil_w   = Inches(2.95)
    pil_h   = Y_MAX - pil_y
    pil_gap = Inches(0.24)

    pillars = [
        ("DATA STRATEGY",       NYU_PURPLE,
         ["▸  Chunked I/O — train.csv streamed, never fully in RAM",
          "▸  Join: building metadata + site weather per chunk",
          "▸  Parquet feature store partitioned by site",
          "▸  Time-based train/val split — no future data leaked"]),
        ("MODELING STRATEGY",   NYU_VIOLET,
         ["▸  Baseline: per-meter historical mean predictor",
          "▸  LightGBM: GOSS + EFB for 50M+ row efficiency",
          "▸  Target: ≥20–40% RMSE improvement over baseline",
          "▸  Per-site models trained in parallel across 19 sites"]),
        ("OUTPUT & EVALUATION", NYU_PURPLE,
         ["▸  Forecast: per-meter hourly predictions (kWh)",
          "▸  Anomaly: Z-score flags + per-meter severity score",
          "▸  Audit list: ranked CSV for operator action",
          "▸  Benchmark table: runtime · memory · speedup"]),
    ]

    for i, (title, col, pts) in enumerate(pillars):
        px = X_L + Inches(i) * (pil_w + pil_gap)
        rect(slide, px, pil_y, pil_w, Inches(0.32), col)
        tb(slide, [WH(title, sz=10)],
           px, pil_y + Inches(0.03), pil_w, Inches(0.28))
        rect(slide, px, pil_y + Inches(0.32), pil_w, pil_h - Inches(0.32), CARD_BG)
        lines = []
        for pt in pts:
            lines += [L(pt, sz=10.5, c=MID), SP(180)]
        tb(slide, lines, px + Inches(0.08), pil_y + Inches(0.38),
           pil_w - Inches(0.12), pil_h - Inches(0.44))


# ── 5. OPTIMIZATION TECHNIQUES ───────────────────────────────────────────────
def s5_optimization(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[7])

    section_title(slide,
                  "OPTIMIZATION TECHNIQUES  &  ADVANCED PYTHON COMPONENTS")

    CONTENT_Y = Y_CONTENT
    COL_H     = Y_MAX - CONTENT_Y
    # 3-column layout
    COL_W  = Inches(2.92)
    COL_GAP= Inches(0.17)

    groups = [
        ("PERFORMANCE & PROFILING", NYU_PURPLE, [
            ("cProfile & memory_profiler",
             "Identify CPU hotspots and peak RAM usage before optimizing. "
             "Benchmark table documents before/after runtime and memory for every component."),
            ("Vectorized NumPy / Pandas",
             "Feature engineering computed with DatetimeIndex, groupby + shift, "
             "and rolling APIs — 100× faster than row-level Python iteration."),
            ("Dtype Optimization",
             "float32 instead of float64 halves memory for 53M-row arrays. "
             "Categorical encoding avoids string overhead in LightGBM features."),
            ("Lazy Iteration with itertools",
             "Meter-site combinations dispatched lazily — no exhaustive "
             "in-memory expansion of the full (3,053 meters × 17,544 hours) index."),
        ]),
        ("PARALLELISM & CONCURRENCY", NYU_VIOLET, [
            ("Multiprocessing",
             "19 sites dispatched to parallel worker processes for independent "
             "feature building and model training. Near-linear speedup measured."),
            ("Concurrent Futures",
             "ThreadPoolExecutor for I/O-bound chunk reads from disk. "
             "ProcessPoolExecutor for CPU-bound feature construction per site."),
            ("Out-of-Core with Dask",
             "Lazy computation graph across 53M rows — same pandas API, "
             "never materializing the full dataset in a single process."),
            ("JIT Compilation — Numba",
             "Modified Z-score anomaly scoring JIT-compiled with @jit(nopython=True). "
             "Target: ≥5× speedup over pure Python loop on 1M+ residuals."),
        ]),
        ("ACCELERATION & BIG DATA", NYU_PURPLE, [
            ("GPU Acceleration — CuPy",
             "CuPy is a drop-in NumPy replacement. Anomaly scoring "
             "and array math run on CUDA without code restructuring."),
            ("LightGBM on GPU",
             "device='gpu' enables GPU-accelerated histogram construction "
             "for GBDT training — benchmarked against CPU baseline."),
            ("Cython Compilation",
             "Custom rolling aggregation kernels compiled to C-speed "
             "when NumPy rolling API is insufficient for the hot path."),
            ("Apache Spark / PySpark",
             "Full 53.6M-row Spark pipeline: distributed feature engineering "
             "via Spark SQL, same CV-RMSE metric at true cluster scale."),
        ]),
    ]

    for col_i, (title, col_h, items) in enumerate(groups):
        cx  = X_L + Inches(col_i) * (COL_W + COL_GAP)
        cy  = CONTENT_Y

        # Column header
        rect(slide, cx, cy, COL_W, Inches(0.36), col_h)
        tb(slide, [WH(title, sz=10.5)],
           cx, cy + Inches(0.04), COL_W, Inches(0.30))

        item_y = cy + Inches(0.40)
        item_h = (COL_H - Inches(0.40)) / len(items)

        for i, (item_title, item_body) in enumerate(items):
            iy = item_y + Inches(i) * item_h
            bg = CARD_BG if i % 2 == 0 else WHITE
            rect(slide, cx, iy, COL_W, item_h - Inches(0.04), bg)
            # Micro accent
            rect(slide, cx, iy, Inches(0.04), item_h - Inches(0.04), col_h)
            tb(slide, [
                L(item_title, sz=11, b=True, c=col_h, f=FONT_H),
                SP(100),
                L(item_body,  sz=10, c=MID),
            ], cx + Inches(0.10), iy + Inches(0.06),
               COL_W - Inches(0.14), item_h - Inches(0.10))


# ── 6. LITERATURE REVIEW ─────────────────────────────────────────────────────
def s6_literature(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[7])

    section_title(slide, "LITERATURE REVIEW")

    papers = [
        ("Miller et al. (2020)",
         "Building Data Genome Project 2  ·  Scientific Data 7, 368",
         "Primary dataset reference. Establishes provenance, quality-control methodology, "
         "and measurement & verification (M&V) use case. Reports gradient boosting "
         "outperforms linear models and neural networks on this data.",
         NYU_PURPLE),
        ("Runge & Zmeureanu (2019)",
         "Forecasting Energy Use in Buildings  ·  Energies 12(18), 3355",
         "Comprehensive ML survey across 140 studies. GBM consistently best on tabular "
         "building data; lag + rolling features are most informative; CV-RMSE is the "
         "recommended cross-building comparison metric.",
         NYU_VIOLET),
        ("Ke et al. (2017)",
         "LightGBM: A Highly Efficient GBDT  ·  NeurIPS 2017",
         "GOSS and EFB optimizations make LightGBM orders of magnitude faster at 50M+ "
         "rows. Native categorical support handles use-type and site features without "
         "one-hot expansion.",
         NYU_PURPLE),
        ("Molina-Solana et al. (2017)",
         "Data Science for Building Energy Mgmt  ·  RSER 70, 598",
         "Widely cited survey establishing residual-based anomaly detection as the "
         "industrially dominant M&V technique (ISO 50001, ASHRAE Guideline 14). "
         "Identifies decision support as the critical gap in academic work.",
         NYU_VIOLET),
        ("Fan et al. (2021)",
         "COVID-19 Impact on Building Energy  ·  Applied Energy 285",
         "Real-world validation: COVID occupancy collapse created massive residuals in "
         "pre-COVID baseline models, correctly flagging abnormal buildings — directly "
         "validating the residual-based anomaly approach.",
         NYU_PURPLE),
        ("Chen & Guestrin (2016)",
         "XGBoost: A Scalable Tree Boosting System  ·  KDD 2016",
         "Baseline comparison model. LightGBM vs. XGBoost at 53M rows provides a "
         "concrete, measurable benchmark demonstrating histogram-based splitting gains "
         "in our pipeline.",
         NYU_VIOLET),
    ]

    CARD_W = Inches(4.52)
    CARD_H = Inches(1.36)
    GAP_X  = Inches(0.22)
    GAP_Y  = Inches(0.10)
    START_Y= Y_CONTENT

    for i, (author, venue, body, col) in enumerate(papers):
        ci = i % 2
        ri = i // 2
        cx = X_L + Inches(ci) * (CARD_W + GAP_X)
        cy = START_Y + Inches(ri) * (CARD_H + GAP_Y)

        rect(slide, cx, cy, CARD_W, CARD_H, CARD_BG)
        rect(slide, cx, cy, Inches(0.055), CARD_H, col)   # left accent bar

        tb(slide, [
            L(author, sz=12, b=True, c=col,     f=FONT_H),
            SP(80),
            L(venue,  sz=9.5, i=True, c="777777"),
            SP(130),
            L(body,   sz=10.5, c=MID),
        ], cx + Inches(0.12), cy + Inches(0.08),
           CARD_W - Inches(0.18), CARD_H - Inches(0.12))

    # Footer note
    footer_y = Y_MAX - Inches(0.28)
    tb(slide, [
        L("Research gap addressed:  Prior work optimizes accuracy OR scales to big data — not both. "
          "Metrik AI bridges this gap with a fully profiled, end-to-end M&V pipeline at 53M-row scale.",
          sz=10.5, b=True, c=NYU_VIOLET, align="ctr"),
    ], X_L, footer_y, SAFE_W, Inches(0.28))


# ── 7. PROJECT GOALS & LONG-TERM VISION ──────────────────────────────────────
def s7_goals(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[7])

    section_title(slide,
                  "PROJECT GOALS  &  LONG-TERM VISION")

    CONTENT_Y = Y_CONTENT
    COL_H     = Y_MAX - CONTENT_Y
    COL_W_L   = Inches(5.55)
    GAP       = Inches(0.18)
    COL_W_R   = SAFE_W - COL_W_L - GAP
    COL_R_X   = X_L + COL_W_L + GAP

    # Vertical divider
    rect(slide, COL_R_X - Inches(0.11), CONTENT_Y, Inches(0.04), COL_H, NYU_VIOLET)

    # ── Left: Project Goals (2 × 2 grid) ─────────────────────────────────────
    goal_w = (COL_W_L - Inches(0.14)) / 2
    goal_h = (COL_H - Inches(0.10)) / 2

    goal_boxes = [
        ("FORECASTING ACCURACY", NYU_PURPLE,
         ["≥20–40% RMSE improvement over per-meter mean baseline",
          "CV-RMSE for fair cross-building comparison",
          "Time-based validation split — no data leakage"]),
        ("ANOMALY DETECTION", NYU_VIOLET,
         ["Residual Z-score flagging per meter",
          "Target ≥5× speedup via JIT compilation",
          "Validated against known fault signatures"]),
        ("SCALABILITY", NYU_PURPLE,
         ["Chunked loader: <4 GB peak for 53.6M rows",
          "Vectorized feature build: ≥2× speedup",
          "≥4-worker parallel training — speedup measured"]),
        ("DECISION SUPPORT OUTPUT", NYU_VIOLET,
         ["Ranked audit list: top-N meters by severity",
          "CLI-generated CSV/JSON — reproducible pipeline",
          "Benchmark table: runtime · memory · speedup"]),
    ]

    for i, (title, col, pts) in enumerate(goal_boxes):
        ci, ri = i % 2, i // 2
        bx = X_L + Inches(ci) * (goal_w + Inches(0.14))
        by = CONTENT_Y + Inches(ri) * (goal_h + Inches(0.10))
        rect(slide, bx, by, goal_w, Inches(0.32), col)
        tb(slide, [WH(title, sz=10)],
           bx, by + Inches(0.03), goal_w, Inches(0.28))
        rect(slide, bx, by + Inches(0.32), goal_w, goal_h - Inches(0.32), CARD_BG)
        lines = []
        for pt in pts:
            lines += [BU(pt, sz=11), SP(200)]
        tb(slide, lines, bx + Inches(0.08), by + Inches(0.38),
           goal_w - Inches(0.12), goal_h - Inches(0.44))

    # ── Right: Long-term Vision ───────────────────────────────────────────────
    tb(slide, [H3("LONG-TERM VISION", sz=11)],
       COL_R_X, CONTENT_Y, COL_W_R, Inches(0.30))

    vision_items = [
        ("Real-Time Deployment",
         "Deploy the pipeline as a live REST API that ingests smart-meter "
         "telemetry and serves per-building hourly forecasts and anomaly "
         "alerts to facility management dashboards."),
        ("Portfolio-Scale Analysis",
         "Scale from 1,636 buildings to university-wide or citywide portfolios "
         "using the PySpark pipeline on a distributed cluster, enabling "
         "cross-building benchmarking and automated retrofitting prioritization."),
        ("Demand Response Integration",
         "Feed forecasts into utility demand response programs — allowing "
         "buildings to automatically pre-condition during off-peak hours and "
         "reduce peak-load charges and grid stress."),
        ("Carbon Accounting",
         "Combine energy forecasts with site-specific grid carbon intensity "
         "data (kg CO₂/kWh) to produce real-time carbon footprint tracking "
         "aligned with net-zero and ESG reporting frameworks."),
    ]

    vy = CONTENT_Y + Inches(0.32)
    item_h = (Y_MAX - vy) / len(vision_items)

    for i, (vt, vb) in enumerate(vision_items):
        iy = vy + Inches(i) * item_h
        rect(slide, COL_R_X, iy, COL_W_R, item_h - Inches(0.06), CARD_BG if i % 2 == 0 else WHITE)
        rect(slide, COL_R_X, iy, Inches(0.04), item_h - Inches(0.06), NYU_VIOLET if i % 2 else NYU_PURPLE)
        tb(slide, [
            L(vt, sz=11, b=True, c=NYU_PURPLE, f=FONT_H),
            SP(100),
            L(vb, sz=10.5, c=MID),
        ], COL_R_X + Inches(0.10), iy + Inches(0.05),
           COL_W_R - Inches(0.14), item_h - Inches(0.10))


# ── 8. THANK YOU ──────────────────────────────────────────────────────────────
def s8_thankyou(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])   # TITLE – purple bg

    set_ph(slide.shapes.placeholders[0], [
        L("THANK YOU", sz=54, b=True, c=WHITE, f=FONT_H),
    ])
    set_ph(slide.shapes.placeholders[1], [
        L("Questions & Discussion", sz=22, c=WHITE, f=FONT_B),
    ])
    try:
        set_ph(slide.shapes.placeholders[2], [
            L("DS-GA 1019  ·  Advanced Python for Data Science  ·  Spring 2026",
              sz=12, c=WHITE, f=FONT_B),
            SP(150),
            L("Dataset: ASHRAE Great Energy Predictor III  ·  Kaggle (2019)",
              sz=11, c=WHITE, f=FONT_B, i=True),
            SP(80),
            L("Reference: Miller et al. (2020), Scientific Data 7, 368",
              sz=11, c=WHITE, f=FONT_B, i=True),
        ])
    except Exception:
        pass

    # Bottom strip
    rect(slide, Inches(0), H - Inches(0.10), W, Inches(0.10), NYU_VIOLET)

    # Summary stats in violet bar above footer
    bar_y = H - Inches(0.88)
    rect(slide, Inches(0), bar_y, W, Inches(0.74), NYU_VIOLET)
    for i, (val, lbl) in enumerate([
        ("53.6M rows",    "ASHRAE dataset"),
        ("3 components",  "Forecast · Anomaly · Decision"),
        ("10 techniques", "Advanced Python tools applied"),
        ("≥5× speedup",   "Numba JIT anomaly scoring target"),
    ]):
        bx = Inches(0.5) + Inches(i * 2.28)
        tb(slide, [WH(val, sz=14), SP(50), WH(lbl, sz=9)],
           bx, bar_y + Inches(0.06), Inches(2.1), Inches(0.65))


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════
def main():
    prs = Presentation(TEMPLATE)
    print("Clearing template slides...")
    clear_slides(prs)

    print("Building slides:")
    s1_title(prs);        print("  [1/8] Title")
    s2_motivation(prs);   print("  [2/8] Motivation + Problem Statement")
    s3_dataset(prs);      print("  [3/8] Dataset")
    s4_pipeline(prs);     print("  [4/8] ML Pipeline (Flowchart)")
    s5_optimization(prs); print("  [5/8] Optimization Techniques")
    s6_literature(prs);   print("  [6/8] Literature Review")
    s7_goals(prs);        print("  [7/8] Goals & Long-term Vision")
    s8_thankyou(prs);     print("  [8/8] Thank You")

    prs.save(OUTPUT)
    print(f"\nSaved → {OUTPUT}")

if __name__ == "__main__":
    main()
