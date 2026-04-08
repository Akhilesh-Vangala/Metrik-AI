"""
Update metrik_ai_v3_final.pptx per user feedback:
1. Add team names on Slide 1
2. Reduce text on Motivation slide to bullets
3. Fix Streamlit wording + add pipeline diagram on Slide 4
4. Compress Literature slide to 3 references
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

path_in = "/Users/akhileshvangala/Downloads/metrik_ai_v3_final.pptx"
path_out = "/Users/akhileshvangala/Desktop/Metrik AI/Proposal Files/metrik_ai_v3_final_updated.pptx"

prs = Presentation(path_in)

# --- SLIDE 1: Add team names under the title area ---
s1 = prs.slides[0]
# Add textbox for team names (below subtitle, centered-ish)
# Typical title slide: title at top, subtitle below. Add names below that.
left = Inches(1)
top = Inches(4.2)
w = Inches(8)
h = Inches(0.8)
tb = s1.shapes.add_textbox(left, top, w, h)
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Akhilesh Vangala (sv3129@nyu.edu)  |  Lucas Yao (ly2808@nyu.edu)"
p.font.size = Pt(14)
p.alignment = PP_ALIGN.CENTER

# --- SLIDE 2: Reduce text to bullets (Shape 1 = long para, Shape 3 = stakeholders) ---
s2 = prs.slides[1]
# Shape 1: Key Problem bullets
tf = s2.shapes[1].text_frame
tf.clear()
p0 = tf.paragraphs[0]
p0.text = "Key Problem"
p0.font.bold = True
p0.font.size = Pt(18)
p0.space_after = Pt(8)
for line in [
    "Buildings consume 40% of global energy",
    "Up to 30% energy waste from inefficient scheduling",
    "$130B annual loss from poor building management",
    "Need predictive, data-driven energy optimization",
]:
    p = tf.add_paragraph()
    p.text = "• " + line
    p.level = 0
    p.font.size = Pt(14)
    p.space_after = Pt(4)
# Shape 3: Who needs this
tf = s2.shapes[3].text_frame
tf.clear()
p0 = tf.paragraphs[0]
p0.text = "Who needs this"
p0.font.bold = True
p0.font.size = Pt(18)
p0.space_after = Pt(8)
for line in [
    "Building operators — forward-looking consumption forecasts",
    "Utilities — demand forecasts for load balancing",
    "ESG & sustainability teams — net-zero benchmarks",
    "Predictive management can cut waste 15–25%, saving $ and CO₂",
]:
    p = tf.add_paragraph()
    p.text = "• " + line
    p.level = 0
    p.font.size = Pt(14)
    p.space_after = Pt(4)

# --- SLIDE 4: Fix Streamlit wording (Shape 3, paragraph index 3) and add pipeline diagram ---
s4 = prs.slides[3]
# Replace Streamlit sentence in shape 3, paragraph 3
s4.shapes[3].text_frame.paragraphs[3].text = (
    "Streamlit dashboard for interactive exploration: forecasts, flagged anomalous meters, and ranked audit list by anomaly severity and savings potential."
)

# Add pipeline diagram as a new text box (vertical flow with arrows)
# Place it on the right or below the main content
left = Inches(0.5)
top = Inches(4.0)
w = Inches(5)
h = Inches(2.2)
tb = s4.shapes.add_textbox(left, top, w, h)
tf = tb.text_frame
tf.word_wrap = False
lines = [
    "ASHRAE Dataset",
    "        ↓",
    "Chunked Data Pipeline",
    "        ↓",
    "Feature Engineering",
    "        ↓",
    "LightGBM Forecast Model",
    "        ↓",
    "Residual Analysis",
    "        ↓",
    "Anomaly Detection",
    "        ↓",
    "Streamlit Dashboard",
]
for i, line in enumerate(lines):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = line
    p.font.size = Pt(12)
    p.font.name = "Consolas"
    p.space_after = Pt(2)

# --- SLIDE 6: Compress Literature to 3 references only (Shape 1 and 3) ---
s6 = prs.slides[5]
tf = s6.shapes[1].text_frame
tf.clear()
p0 = tf.paragraphs[0]
p0.text = "Key references"
p0.font.bold = True
p0.font.size = Pt(18)
p0.space_after = Pt(10)
for r in [
    "Miller et al. (2020) — Building Data Genome dataset",
    "Ke et al. (2017) — LightGBM architecture",
    "Molina-Solana et al. (2017) — ML for energy management",
]:
    p = tf.add_paragraph()
    p.text = "• " + r
    p.level = 0
    p.font.size = Pt(14)
    p.space_after = Pt(6)
tf = s6.shapes[3].text_frame
tf.clear()
p = tf.paragraphs[0]
p.text = "(These three references validate our dataset, model choice, and M&V workflow.)"
p.font.size = Pt(11)
try:
    p.font.italic = True
except Exception:
    pass

prs.save(path_out)
print("Saved to:", path_out)
