"""
build_ppt.py — Generates the DS-GA 1019 Proposal Presentation for Metrik AI
Run: python3 build_ppt.py
Output: MetrikAI_Proposal_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color palette ────────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
BG_CARD      = RGBColor(0x15, 0x2A, 0x3E)   # slightly lighter navy for cards
ACCENT_TEAL  = RGBColor(0x00, 0xC2, 0xB5)   # teal / cyan accent
ACCENT_GOLD  = RGBColor(0xF5, 0xA6, 0x23)   # warm gold
ACCENT_CORAL = RGBColor(0xFF, 0x6B, 0x6B)   # soft coral for highlights
TEXT_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT   = RGBColor(0xC8, 0xD6, 0xE5)   # muted blue-white for body
TEXT_DIM     = RGBColor(0x7F, 0x9A, 0xB3)   # dim for captions
DIVIDER      = RGBColor(0x00, 0xC2, 0xB5)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ── Helpers ───────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_layout(prs):
    return prs.slide_layouts[6]   # completely blank


def add_bg(slide, color=BG_DARK):
    """Fill slide background with solid color."""
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()   # no border
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=TEXT_WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def add_multiline_textbox(slide, lines, left, top, width, height,
                          default_size=16, default_color=TEXT_LIGHT,
                          default_bold=False, line_spacing=None):
    """lines = list of (text, size, color, bold, italic, align) or just strings."""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for item in lines:
        if isinstance(item, str):
            txt, sz, col, bld, itl, aln = item, default_size, default_color, default_bold, False, PP_ALIGN.LEFT
        else:
            txt = item.get("text", "")
            sz  = item.get("size", default_size)
            col = item.get("color", default_color)
            bld = item.get("bold", default_bold)
            itl = item.get("italic", False)
            aln = item.get("align", PP_ALIGN.LEFT)

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.alignment = aln
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(sz)
        run.font.bold = bld
        run.font.italic = itl
        run.font.color.rgb = col
        run.font.name = "Calibri"
    return txBox


def divider_line(slide, y, left=Inches(0.5), width=Inches(12.33), color=DIVIDER, thickness=Pt(1.5)):
    from pptx.util import Pt as PtU
    line = slide.shapes.add_shape(1, left, y, width, Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def slide_number_tag(slide, num, total=11):
    add_textbox(slide, f"{num} / {total}",
                left=Inches(12.2), top=Inches(7.1), width=Inches(1.0), height=Inches(0.3),
                font_size=10, color=TEXT_DIM, align=PP_ALIGN.RIGHT)


def section_label(slide, label, left=Inches(0.5), top=Inches(0.18)):
    add_textbox(slide, label.upper(),
                left=left, top=top, width=Inches(4), height=Inches(0.3),
                font_size=9, color=ACCENT_TEAL, bold=True)


def slide_title(slide, title, subtitle=None, top=Inches(0.55)):
    add_textbox(slide, title,
                left=Inches(0.5), top=top, width=Inches(12.33), height=Inches(0.7),
                font_size=32, bold=True, color=TEXT_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        divider_line(slide, top + Inches(0.72))
        add_textbox(slide, subtitle,
                    left=Inches(0.5), top=top + Inches(0.78), width=Inches(12.33), height=Inches(0.35),
                    font_size=14, color=ACCENT_TEAL, bold=False, align=PP_ALIGN.LEFT)


def bullet_block(slide, items, left, top, width, height,
                 bullet="▸", size=15, color=TEXT_LIGHT, head_color=ACCENT_TEAL,
                 head_size=16, gap=True):
    """items = list of strings or (header, body) tuples."""
    lines = []
    for item in items:
        if isinstance(item, tuple):
            h, b = item
            lines.append({"text": h, "size": head_size, "color": head_color, "bold": True})
            if b:
                lines.append({"text": f"  {b}", "size": size, "color": color, "bold": False})
        else:
            lines.append({"text": f"{bullet}  {item}", "size": size, "color": color, "bold": False})
        if gap:
            lines.append({"text": "", "size": 6, "color": color})
    add_multiline_textbox(slide, lines, left, top, width, height)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Title Slide"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide, BG_DARK)

    # Left accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT_TEAL)

    # Top decorative stripe
    add_rect(slide, Inches(0.08), Inches(0), SLIDE_W, Inches(0.04), ACCENT_GOLD)

    # Course label top-right
    add_textbox(slide, "DS-GA 1019  ·  Advanced Python for Data Science  ·  Spring 2026",
                left=Inches(0.5), top=Inches(0.25), width=Inches(12), height=Inches(0.3),
                font_size=11, color=TEXT_DIM, align=PP_ALIGN.LEFT)

    # Main title
    add_textbox(slide, "Metrik AI",
                left=Inches(0.5), top=Inches(1.8), width=Inches(12), height=Inches(1.2),
                font_size=64, bold=True, color=TEXT_WHITE, align=PP_ALIGN.LEFT)

    # Teal accent underline
    add_rect(slide, Inches(0.5), Inches(3.05), Inches(5.5), Inches(0.06), ACCENT_TEAL)

    # Subtitle
    add_textbox(slide, "Predicting & Optimizing Building Energy Consumption at Scale",
                left=Inches(0.5), top=Inches(3.2), width=Inches(12), height=Inches(0.6),
                font_size=22, bold=False, color=ACCENT_TEAL, align=PP_ALIGN.LEFT)

    # Description line
    add_textbox(slide,
                "An end-to-end advanced Python pipeline: forecasting · anomaly detection · decision support",
                left=Inches(0.5), top=Inches(3.85), width=Inches(12), height=Inches(0.4),
                font_size=14, color=TEXT_LIGHT, align=PP_ALIGN.LEFT)

    # Bottom meta
    add_rect(slide, Inches(0.5), Inches(6.3), Inches(12.33), Inches(0.003), ACCENT_TEAL)
    add_textbox(slide, "Project Proposal  ·  March 4, 2026  ·  NYU Center for Data Science",
                left=Inches(0.5), top=Inches(6.4), width=Inches(8), height=Inches(0.3),
                font_size=11, color=TEXT_DIM)

    # Stat boxes bottom-right
    for i, (val, lbl) in enumerate([("53.6M", "meter readings"), ("1,636", "buildings"), ("19", "sites")]):
        bx = Inches(9.0) + Inches(i * 1.45)
        add_rect(slide, bx, Inches(6.1), Inches(1.3), Inches(0.9), BG_CARD)
        add_textbox(slide, val, left=bx, top=Inches(6.12), width=Inches(1.3), height=Inches(0.45),
                    font_size=20, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
        add_textbox(slide, lbl, left=bx, top=Inches(6.56), width=Inches(1.3), height=Inches(0.25),
                    font_size=9, color=TEXT_DIM, align=PP_ALIGN.CENTER)

    slide_number_tag(slide, 1)
    return slide


def slide_02_motivation(prs):
    """Real-World Motivation"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT_TEAL)

    section_label(slide, "Real-World Motivation")
    slide_title(slide, "Why Does This Problem Exist?",
                subtitle="Buildings are the largest, most solvable energy challenge on the planet.")

    # 4 stat cards top row
    stats = [
        ("40%",     "of U.S. energy\nconsumed by buildings",         ACCENT_TEAL),
        ("28%",     "of global CO₂\nfrom building operations",       ACCENT_GOLD),
        ("20–30%",  "of building energy\ntypically wasted",          ACCENT_CORAL),
        ("$130B",   "annual cost of\nbuilding energy waste",         RGBColor(0xA8, 0x8E, 0xFF)),
    ]
    for i, (val, lbl, col) in enumerate(stats):
        bx = Inches(0.5) + Inches(i * 3.1)
        add_rect(slide, bx, Inches(1.65), Inches(2.85), Inches(1.35), BG_CARD)
        add_rect(slide, bx, Inches(1.65), Inches(2.85), Inches(0.06), col)
        add_textbox(slide, val, left=bx, top=Inches(1.72), width=Inches(2.85), height=Inches(0.6),
                    font_size=30, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_textbox(slide, lbl, left=bx, top=Inches(2.32), width=Inches(2.85), height=Inches(0.55),
                    font_size=12, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    # Stakeholders
    divider_line(slide, Inches(3.2))
    add_textbox(slide, "WHO HAS THIS PROBLEM?",
                left=Inches(0.5), top=Inches(3.28), width=Inches(5), height=Inches(0.3),
                font_size=10, color=ACCENT_TEAL, bold=True)
    stakeholders = [
        "Building operators & facility managers — lack forward-looking consumption data",
        "Utilities & grid operators — need demand forecasts for load balancing",
        "ESG & sustainability teams — need quantifiable benchmarks for net-zero targets",
        "Retrofitting decision-makers — need to prioritize which buildings to upgrade",
    ]
    bullet_block(slide, stakeholders,
                 left=Inches(0.5), top=Inches(3.62), width=Inches(6.1), height=Inches(2.8),
                 size=14, gap=False)

    # Core tension box
    add_rect(slide, Inches(6.9), Inches(3.2), Inches(5.9), Inches(3.2), BG_CARD)
    add_rect(slide, Inches(6.9), Inches(3.2), Inches(5.9), Inches(0.06), ACCENT_GOLD)
    add_textbox(slide, "THE CORE TENSION",
                left=Inches(7.0), top=Inches(3.3), width=Inches(5.7), height=Inches(0.3),
                font_size=10, bold=True, color=ACCENT_GOLD)
    add_textbox(slide,
                "Operators must choose between over-heating / cooling (waste + cost) or under-delivering "
                "(occupant discomfort). Without accurate forecasts, neither side of this trade-off is winnable.",
                left=Inches(7.0), top=Inches(3.65), width=Inches(5.7), height=Inches(1.1),
                font_size=13, color=TEXT_LIGHT, italic=True)

    add_textbox(slide, "WHY NOW?",
                left=Inches(7.0), top=Inches(4.85), width=Inches(5.7), height=Inches(0.3),
                font_size=10, bold=True, color=ACCENT_TEAL)
    add_textbox(slide,
                "Smart meters + IoT + open datasets like ASHRAE now make predictive, "
                "data-driven energy management possible — but only if you can process "
                "50M+ records efficiently. That's exactly what this course teaches.",
                left=Inches(7.0), top=Inches(5.2), width=Inches(5.7), height=Inches(1.0),
                font_size=13, color=TEXT_LIGHT)

    slide_number_tag(slide, 2)
    return slide


def slide_03_problem(prs):
    """Problem Statement"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT_CORAL)

    section_label(slide, "Problem Statement", top=Inches(0.18))
    slide_title(slide, "Formal Problem Definition",
                subtitle="Prediction · Anomaly Detection · Decision Support — all at scale.")

    # Formal statement box
    add_rect(slide, Inches(0.5), Inches(1.55), Inches(12.33), Inches(1.1), BG_CARD)
    add_rect(slide, Inches(0.5), Inches(1.55), Inches(0.06), Inches(1.1), ACCENT_CORAL)
    add_textbox(slide,
                "Given hourly meter readings from 1,636 non-residential buildings across 19 sites over two "
                "full years, along with building metadata and weather data, predict next-hour energy "
                "consumption at the individual meter level — and flag anomalies and prioritize interventions "
                "— without exceeding commodity memory limits.",
                left=Inches(0.7), top=Inches(1.62), width=Inches(11.9), height=Inches(0.9),
                font_size=14, color=TEXT_WHITE, italic=True)

    # Three components
    divider_line(slide, Inches(2.8))
    components = [
        ("FORECASTING",      ACCENT_TEAL,  "Predict next-hour meter-level consumption\n→ Enables pre-cooling/heating & demand response"),
        ("ANOMALY DETECTION", ACCENT_GOLD, "Flag residuals = actual − predicted > threshold\n→ Catches meter faults, equipment left on, billing errors"),
        ("DECISION SUPPORT",  ACCENT_CORAL,"Rank buildings/meters by anomaly severity\n→ Tells operators where to act first"),
    ]
    for i, (title, col, body) in enumerate(components):
        bx = Inches(0.5) + Inches(i * 4.1)
        add_rect(slide, bx, Inches(2.9), Inches(3.9), Inches(1.65), col)
        add_textbox(slide, title, left=bx, top=Inches(2.92), width=Inches(3.9), height=Inches(0.45),
                    font_size=13, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        add_rect(slide, bx, Inches(3.6), Inches(3.9), Inches(1.55), BG_CARD)
        add_textbox(slide, body, left=bx + Inches(0.12), top=Inches(3.65), width=Inches(3.66), height=Inches(1.4),
                    font_size=13, color=TEXT_LIGHT)

    # Why hard
    add_textbox(slide, "WHAT MAKES THIS HARD",
                left=Inches(0.5), top=Inches(5.3), width=Inches(6), height=Inches(0.28),
                font_size=10, color=ACCENT_TEAL, bold=True)
    hard = [
        "Scale: 53.6M rows — cannot naively load into RAM",
        "Temporal leakage risk: rolling/lag features must be strictly past-only",
        "Heterogeneity: 4 meter types · 16 building categories · 19 sites · mixed units",
        "Missing data: weather gaps, meter outages, negative readings",
    ]
    bullet_block(slide, hard, left=Inches(0.5), top=Inches(5.6), width=Inches(12.33), height=Inches(1.5),
                 size=13, gap=False)

    slide_number_tag(slide, 3)
    return slide


def slide_04_solution(prs):
    """Proposed Solution"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT_GOLD)

    section_label(slide, "Proposed Solution")
    slide_title(slide, "What We Build & What We Deliver",
                subtitle="Three-component pipeline · Out-of-core · Fully profiled & benchmarked")

    # Pipeline flow diagram (text-based)
    steps = [
        ("train.csv\n53.6M rows",        ACCENT_TEAL),
        ("Chunked\nReader",               BG_CARD),
        ("Feature\nEngineering",          BG_CARD),
        ("LightGBM\nForecaster",          BG_CARD),
        ("Numba\nAnomaly Score",          BG_CARD),
        ("Decision\nSupport Output",      ACCENT_GOLD),
    ]
    bw, bh = Inches(1.75), Inches(0.75)
    for i, (label, col) in enumerate(steps):
        bx = Inches(0.38) + Inches(i * 2.1)
        by = Inches(1.65)
        add_rect(slide, bx, by, bw, bh, col)
        txt_col = BG_DARK if col in (ACCENT_TEAL, ACCENT_GOLD) else TEXT_LIGHT
        add_textbox(slide, label, left=bx, top=by, width=bw, height=bh,
                    font_size=11, bold=True, color=txt_col, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_textbox(slide, "→",
                        left=bx + bw, top=by + Inches(0.18), width=Inches(0.35), height=Inches(0.35),
                        font_size=18, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

    # Three component detail boxes
    cols_data = [
        ("FORECASTING",      ACCENT_TEAL,
         ["Model: LightGBM vs. per-meter mean baseline",
          "Features: lags, rolling means, time, weather",
          "Split: time-based (last 3 months = validation)",
          "Metric: RMSE · CV-RMSE per meter",
          "Target: ≥20–40% RMSE gain over baseline"]),
        ("ANOMALY DETECTION", ACCENT_GOLD,
         ["Input: residuals = actual − predicted",
          "Method: Modified Z-score / MAD scoring",
          "Output: per-meter anomaly score + flag",
          "Engine: Numba JIT over 1M+ residuals",
          "Target: ≥5× speedup vs. Python loop"]),
        ("DECISION SUPPORT",  ACCENT_CORAL,
         ["Input: forecasts + anomaly scores",
          "Output: ranked audit list (top-N meters)",
          "Sort by: anomaly severity, excess kWh",
          "Delivery: CLI-generated CSV/JSON report",
          "Mirrors real ASHRAE Guideline 14 M&V workflow"]),
    ]
    for i, (title, col, pts) in enumerate(cols_data):
        bx = Inches(0.5) + Inches(i * 4.28)
        add_rect(slide, bx, Inches(2.65), Inches(4.0), Inches(0.35), col)
        add_textbox(slide, title, left=bx, top=Inches(2.67), width=Inches(4.0), height=Inches(0.32),
                    font_size=11, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        add_rect(slide, bx, Inches(3.0), Inches(4.0), Inches(2.1), BG_CARD)
        for j, pt in enumerate(pts):
            add_textbox(slide, f"▸  {pt}",
                        left=bx + Inches(0.12), top=Inches(3.08) + Inches(j * 0.39),
                        width=Inches(3.75), height=Inches(0.37),
                        font_size=12, color=TEXT_LIGHT)

    # Engineering strip
    divider_line(slide, Inches(5.3))
    eng_items = ["src/ · scripts/ · config/ · tests/ package structure",
                 "CLI: python -m src.cli run --config config.yaml",
                 "requirements.txt · pinned seeds · README",
                 "Benchmark table: before/after runtime + speedup"]
    add_textbox(slide, "ENGINEERING DELIVERABLES",
                left=Inches(0.5), top=Inches(5.35), width=Inches(4), height=Inches(0.28),
                font_size=10, color=ACCENT_TEAL, bold=True)
    for i, item in enumerate(eng_items):
        col_pos = Inches(0.5) + Inches(i * 3.1)
        add_rect(slide, col_pos, Inches(5.68), Inches(2.95), Inches(0.55), BG_CARD)
        add_textbox(slide, item, left=col_pos + Inches(0.1), top=Inches(5.71),
                    width=Inches(2.75), height=Inches(0.5), font_size=11, color=TEXT_LIGHT)

    slide_number_tag(slide, 4)
    return slide


def slide_05_dataset_overview(prs):
    """Dataset Overview"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT_TEAL)

    section_label(slide, "Dataset")
    slide_title(slide, "ASHRAE Great Energy Predictor III",
                subtitle="The largest publicly available benchmark dataset for building energy prediction")

    # Left panel: scale stats
    scale_stats = [
        ("53.6M",   "hourly training rows",        ACCENT_TEAL),
        ("1,636",   "non-residential buildings",    ACCENT_GOLD),
        ("3,053",   "energy meters",                ACCENT_CORAL),
        ("19",      "geographic sites (NA + EU)",   RGBColor(0xA8, 0x8E, 0xFF)),
        ("2 years", "Jan 2016 – Dec 2017",          TEXT_LIGHT),
        ("Hourly",  "17,544 hrs/meter over 2 yrs",  TEXT_LIGHT),
    ]
    for i, (val, lbl, col) in enumerate(scale_stats):
        by = Inches(1.65) + Inches(i * 0.87)
        add_rect(slide, Inches(0.5), by, Inches(4.1), Inches(0.82), BG_CARD)
        add_textbox(slide, val, left=Inches(0.62), top=by + Inches(0.04),
                    width=Inches(1.6), height=Inches(0.42),
                    font_size=22, bold=True, color=col, align=PP_ALIGN.LEFT)
        add_textbox(slide, lbl, left=Inches(2.25), top=by + Inches(0.15),
                    width=Inches(2.2), height=Inches(0.42),
                    font_size=12, color=TEXT_LIGHT, align=PP_ALIGN.LEFT)

    # Right panel: provenance & significance
    add_rect(slide, Inches(5.0), Inches(1.55), Inches(7.8), Inches(2.4), BG_CARD)
    add_rect(slide, Inches(5.0), Inches(1.55), Inches(7.8), Inches(0.05), ACCENT_TEAL)
    add_textbox(slide, "DATASET PROVENANCE & SIGNIFICANCE",
                left=Inches(5.1), top=Inches(1.62), width=Inches(7.5), height=Inches(0.3),
                font_size=10, bold=True, color=ACCENT_TEAL)
    prov = [
        "Kaggle competition hosted by ASHRAE (Oct–Dec 2019) · 4,370 competing teams",
        "Grounded in the Building Data Genome Project 2 (BDG2) — peer-reviewed, real operational data",
        "Published: Miller et al. (2020), Scientific Data 7, 368 · arXiv:2006.02273",
        "Cited in 100+ academic papers on building energy, transfer learning, demand response",
        "Benchmark dataset for IEEE Power & Energy Society and DOE building research",
        "Non-residential only: offices, retail, education, lodging, healthcare, warehouses",
    ]
    for i, pt in enumerate(prov):
        add_textbox(slide, f"▸  {pt}",
                    left=Inches(5.1), top=Inches(2.0) + Inches(i * 0.32),
                    width=Inches(7.6), height=Inches(0.3),
                    font_size=12, color=TEXT_LIGHT)

    # Files table
    add_textbox(slide, "FILES INCLUDED",
                left=Inches(5.0), top=Inches(4.1), width=Inches(4), height=Inches(0.28),
                font_size=10, bold=True, color=ACCENT_TEAL)
    files = [
        ("train.csv",              "53.6M rows · building, meter, timestamp, reading (target)"),
        ("building_metadata.csv",  "1,636 rows · site, primary use, sq footage, year built"),
        ("weather_train.csv",      "~332K rows · hourly weather per site (temp, dew, wind…)"),
        ("test.csv",               "Competition test set · same schema, no target"),
        ("sample_submission.csv",  "Submission format: row_id + predicted meter_reading"),
    ]
    for i, (fname, desc) in enumerate(files):
        by = Inches(4.45) + Inches(i * 0.54)
        add_rect(slide, Inches(5.0), by, Inches(7.8), Inches(0.5), BG_CARD if i % 2 == 0 else BG_DARK)
        add_textbox(slide, fname, left=Inches(5.1), top=by + Inches(0.06),
                    width=Inches(2.2), height=Inches(0.38),
                    font_size=12, bold=True, color=ACCENT_GOLD)
        add_textbox(slide, desc, left=Inches(7.35), top=by + Inches(0.06),
                    width=Inches(5.3), height=Inches(0.38),
                    font_size=12, color=TEXT_LIGHT)

    slide_number_tag(slide, 5)
    return slide


def slide_06_dataset_deep(prs):
    """Dataset Deep Dive — Schema & Features"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT_GOLD)

    section_label(slide, "Dataset — Deep Dive")
    slide_title(slide, "Schema, Features & Data Characteristics",
                subtitle="Target variable · Feature engineering strategy · Modeling-relevant properties")

    # Target variable box
    add_rect(slide, Inches(0.5), Inches(1.55), Inches(12.33), Inches(0.75), BG_CARD)
    add_rect(slide, Inches(0.5), Inches(1.55), Inches(0.06), Inches(0.75), ACCENT_CORAL)
    add_textbox(slide, "TARGET:  meter_reading  (float, kWh or kBTU)  —  hourly energy consumption per meter  —  continuous regression",
                left=Inches(0.7), top=Inches(1.63), width=Inches(11.9), height=Inches(0.55),
                font_size=13, bold=False, color=TEXT_WHITE)

    # Meter types
    add_textbox(slide, "METER TYPES",
                left=Inches(0.5), top=Inches(2.42), width=Inches(4), height=Inches(0.28),
                font_size=10, bold=True, color=ACCENT_TEAL)
    meters = [("0 — Electricity", "Whole-building · kWh (kBTU for Site 0)"),
              ("1 — Chilled Water", "Cooling energy"),
              ("2 — Steam", "Heating energy"),
              ("3 — Hot Water", "Heating energy")]
    for i, (code, desc) in enumerate(meters):
        add_textbox(slide, f"{code}  ·  {desc}",
                    left=Inches(0.5), top=Inches(2.75) + Inches(i * 0.32),
                    width=Inches(5.5), height=Inches(0.3),
                    font_size=12, color=TEXT_LIGHT)

    # Building metadata columns
    add_textbox(slide, "BUILDING METADATA FEATURES",
                left=Inches(0.5), top=Inches(4.1), width=Inches(5), height=Inches(0.28),
                font_size=10, bold=True, color=ACCENT_TEAL)
    bm_cols = [
        "site_id  (0–18)  ·  19 geographic clusters",
        "primary_use  ·  16 categories (Education, Office, Lodging, Retail, Healthcare…)",
        "square_feet  ·  ~300 to ~875,000 sq ft",
        "year_built  ·  ~50% missing — requires imputation strategy",
        "floor_count  ·  highly sparse",
    ]
    for i, col in enumerate(bm_cols):
        add_textbox(slide, f"▸  {col}",
                    left=Inches(0.5), top=Inches(4.45) + Inches(i * 0.32),
                    width=Inches(6.0), height=Inches(0.3),
                    font_size=12, color=TEXT_LIGHT)

    # Engineered features table
    add_textbox(slide, "ENGINEERED FEATURES",
                left=Inches(6.9), top=Inches(2.42), width=Inches(5), height=Inches(0.28),
                font_size=10, bold=True, color=ACCENT_TEAL)
    feat_rows = [
        ("Temporal",  "hour, dow, month, is_weekend, is_holiday  ·  vectorized via DatetimeIndex"),
        ("Lag",       "same_hour_yesterday (lag=24), same_day_last_week (lag=168)  ·  groupby + shift"),
        ("Rolling",   "rolling_mean_24h, rolling_mean_168h  ·  past-only window  ·  leakage-safe"),
        ("Building",  "log_sq_ft, building_age, primary_use_encoded  ·  joined at load time"),
        ("Weather",   "air_temperature, rolling_temp_24h, temp_lag_1  ·  per (site_id, timestamp)"),
    ]
    for i, (grp, desc) in enumerate(feat_rows):
        by = Inches(2.78) + Inches(i * 0.58)
        add_rect(slide, Inches(6.9), by, Inches(6.0), Inches(0.54), BG_CARD if i % 2 == 0 else BG_DARK)
        add_textbox(slide, grp, left=Inches(7.0), top=by + Inches(0.05),
                    width=Inches(1.1), height=Inches(0.44),
                    font_size=12, bold=True, color=ACCENT_GOLD)
        add_textbox(slide, desc, left=Inches(8.15), top=by + Inches(0.05),
                    width=Inches(4.6), height=Inches(0.44),
                    font_size=11, color=TEXT_LIGHT)

    # Data characteristics
    add_textbox(slide, "KEY DATA CHARACTERISTICS",
                left=Inches(6.9), top=Inches(5.7), width=Inches(5), height=Inches(0.28),
                font_size=10, bold=True, color=ACCENT_GOLD)
    chars = ["Strong daily & weekly seasonality  ·  Hierarchical: Site → Building → Meter",
             "Weather-driven HVAC load  ·  Site-level weather, meter-level consumption",
             "Heterogeneous: campus ≠ hotel ≠ parking garage  ·  models must handle this"]
    for i, c in enumerate(chars):
        add_textbox(slide, f"▸  {c}",
                    left=Inches(6.9), top=Inches(6.05) + Inches(i * 0.35),
                    width=Inches(6.0), height=Inches(0.32),
                    font_size=12, color=TEXT_LIGHT)

    slide_number_tag(slide, 6)
    return slide


def slide_07_technical_approach(prs):
    """Technical Approach & Benchmark Targets"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT_TEAL)

    section_label(slide, "Technical Approach")
    slide_title(slide, "End-to-End Methodology & Optimization Strategy",
                subtitle="Out-of-core loading · Vectorized features · LightGBM · Numba · Parallel training · PySpark")

    # 5 pipeline steps
    steps = [
        ("1. Out-of-Core Load",       ACCENT_TEAL,
         "pd.read_csv(chunksize=2M) or Dask\nJoin metadata + weather per chunk\nTarget: <4 GB peak memory"),
        ("2. Feature Engineering",    ACCENT_GOLD,
         "100% vectorized (no .iterrows)\nLag + rolling + time + building + weather\nUnit tests: zero future leakage"),
        ("3. LightGBM Forecasting",   ACCENT_CORAL,
         "Baseline: per-meter historical mean\nModel: LightGBM (GOSS + EFB)\nTarget: ≥20–40% RMSE gain"),
        ("4. Anomaly Scoring",        RGBColor(0xA8, 0x8E, 0xFF),
         "Residuals: actual − predicted\nModified Z-score / MAD per meter\nNumba JIT · target ≥5× speedup"),
        ("5. Decision Support",       RGBColor(0x4E, 0xC9, 0x71),
         "Rank by anomaly score + excess kWh\nCLI output: audit_list.csv\nMirrors ASHRAE Guideline 14 M&V"),
    ]
    for i, (title, col, body) in enumerate(steps):
        bx = Inches(0.38) + Inches(i * 2.55)
        add_rect(slide, bx, Inches(1.65), Inches(2.4), Inches(0.38), col)
        add_textbox(slide, title, left=bx, top=Inches(1.67), width=Inches(2.4), height=Inches(0.35),
                    font_size=11, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        add_rect(slide, bx, Inches(2.03), Inches(2.4), Inches(1.65), BG_CARD)
        add_textbox(slide, body, left=bx + Inches(0.1), top=Inches(2.1),
                    width=Inches(2.2), height=Inches(1.5),
                    font_size=11, color=TEXT_LIGHT)

    # Benchmark table
    divider_line(slide, Inches(3.85))
    add_textbox(slide, "BENCHMARK TARGETS",
                left=Inches(0.5), top=Inches(3.92), width=Inches(5), height=Inches(0.28),
                font_size=10, bold=True, color=ACCENT_TEAL)

    headers = ["Component", "Baseline Method", "Optimized Method", "Success Target"]
    col_widths = [Inches(2.4), Inches(2.8), Inches(3.3), Inches(4.1)]
    col_starts = [Inches(0.5), Inches(2.95), Inches(5.8), Inches(9.15)]

    # Header row
    for j, (hdr, cw, cs) in enumerate(zip(headers, col_widths, col_starts)):
        add_rect(slide, cs, Inches(4.28), cw, Inches(0.35), ACCENT_TEAL)
        add_textbox(slide, hdr, left=cs + Inches(0.05), top=Inches(4.3),
                    width=cw - Inches(0.1), height=Inches(0.3),
                    font_size=11, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

    bm_rows = [
        ("Data Load",          "Full in-memory (OOM)",    "Chunked / Dask out-of-core",   "No OOM · <4 GB peak"),
        ("Feature Build (5M)", "Single-thread loop",      "Vectorized pandas + parallel",  "≥2× speedup"),
        ("Forecasting",        "Per-meter mean",          "LightGBM (GOSS + EFB)",         "≥20–40% RMSE gain"),
        ("Anomaly (1M rows)",  "Pure Python loop",        "Numba @jit(nopython=True)",     "≥5× speedup"),
        ("End-to-End (5M)",    "Sequential pipeline",     "Parallel + profiled pipeline",  "Runtime <10 min"),
    ]
    for i, row in enumerate(bm_rows):
        by = Inches(4.63) + Inches(i * 0.42)
        row_bg = BG_CARD if i % 2 == 0 else BG_DARK
        for j, (cell, cw, cs) in enumerate(zip(row, col_widths, col_starts)):
            add_rect(slide, cs, by, cw, Inches(0.38), row_bg)
            cell_col = ACCENT_GOLD if j == 0 else (RGBColor(0x4E, 0xC9, 0x71) if j == 3 else TEXT_LIGHT)
            add_textbox(slide, cell, left=cs + Inches(0.05), top=by + Inches(0.04),
                        width=cw - Inches(0.1), height=Inches(0.32),
                        font_size=11, bold=(j == 0), color=cell_col)

    slide_number_tag(slide, 7)
    return slide


def slide_08_course_techniques(prs):
    """Course Techniques — full mapping"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT_GOLD)

    section_label(slide, "Course Techniques — DS-GA 1019")
    slide_title(slide, "Every Course Topic Has a Job in This Pipeline",
                subtitle="10 of 13 course weeks directly applied — not demonstrations, but load-bearing components")

    # Table of course mappings
    course_items = [
        ("Wk 2", "Python Performance Tips",  "cProfile + memory_profiler; dtype optimization (float32 → halves 53M-row memory)", ACCENT_TEAL),
        ("Wk 3", "itertools Module",          "Lazy dispatch of (meter, site) pairs; itertools.chain for streaming site partitions", ACCENT_TEAL),
        ("Wk 4", "Performance Tuning",        "Vectorized feature build: DatetimeIndex.hour is 100× faster than apply(lambda x: x.hour)", ACCENT_GOLD),
        ("Wk 5", "Cython",                    "Custom rolling window kernels compiled to C when pandas rolling API is insufficient", ACCENT_GOLD),
        ("Wk 6", "Numba",                     "Anomaly scoring: @jit(nopython=True) over 1M+ residuals → target ≥5× speedup", ACCENT_CORAL),
        ("Wk 8", "Optimization in Python",    "Dict-based weather join O(1) vs. DataFrame filter O(n); Bayesian hyperparameter search", ACCENT_CORAL),
        ("Wk 9", "Python Concurrency",        "ThreadPoolExecutor (I/O-bound chunk reads) vs. ProcessPoolExecutor (CPU-bound feature build)", RGBColor(0xA8, 0x8E, 0xFF)),
        ("Wk 10–11", "Parallel Programming",  "Embarrassingly parallel per-site training: 19 sites → 19 workers → near-linear speedup", RGBColor(0xA8, 0x8E, 0xFF)),
        ("Wk 12", "Python for GPUs",          "CuPy drop-in NumPy replacement for anomaly scoring; LightGBM device='gpu' training", RGBColor(0x4E, 0xC9, 0x71)),
        ("Wk 13", "BigData with PySpark",     "Full 53M-row Spark pipeline: spark.read.csv → Spark SQL features → same CV-RMSE metric", RGBColor(0x4E, 0xC9, 0x71)),
    ]

    col1_w = Inches(0.65)   # week
    col2_w = Inches(1.95)   # topic
    col3_w = Inches(9.45)   # role
    col1_x = Inches(0.4)
    col2_x = Inches(1.1)
    col3_x = Inches(3.1)
    row_h  = Inches(0.5)

    # Header
    for cx, cw, hdr in [(col1_x, col1_w, "Week"), (col2_x, col2_w, "Topic"), (col3_x, col3_w, "Role in Metrik AI")]:
        add_rect(slide, cx, Inches(1.58), cw, Inches(0.32), ACCENT_TEAL)
        add_textbox(slide, hdr, left=cx + Inches(0.04), top=Inches(1.6),
                    width=cw - Inches(0.08), height=Inches(0.28),
                    font_size=11, bold=True, color=BG_DARK)

    for i, (wk, topic, role, col) in enumerate(course_items):
        by = Inches(1.9) + Inches(i * row_h)
        row_bg = BG_CARD if i % 2 == 0 else BG_DARK
        add_rect(slide, col1_x, by, col1_w, row_h - Inches(0.02), row_bg)
        add_rect(slide, col2_x, by, col2_w, row_h - Inches(0.02), row_bg)
        add_rect(slide, col3_x, by, col3_w, row_h - Inches(0.02), row_bg)

        add_textbox(slide, wk, left=col1_x + Inches(0.04), top=by + Inches(0.09),
                    width=col1_w, height=row_h,
                    font_size=10, bold=True, color=col)
        add_textbox(slide, topic, left=col2_x + Inches(0.04), top=by + Inches(0.09),
                    width=col2_w, height=row_h,
                    font_size=11, bold=True, color=TEXT_WHITE)
        add_textbox(slide, role, left=col3_x + Inches(0.06), top=by + Inches(0.09),
                    width=col3_w - Inches(0.1), height=row_h,
                    font_size=11, color=TEXT_LIGHT)

    slide_number_tag(slide, 8)
    return slide


def slide_09_literature(prs):
    """Literature Review"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT_CORAL)

    section_label(slide, "Literature Review")
    slide_title(slide, "Prior Work & How We Extend It",
                subtitle="Grounded in peer-reviewed research · We address the scalability + end-to-end gap")

    papers = [
        ("Miller et al. (2020)",
         "Building Data Genome Project 2  ·  Scientific Data 7, 368",
         "Primary dataset reference. Establishes BDG2 provenance, quality checks, and M&V use case. "
         "Reports gradient boosting outperforms linear models and NNs on this data → validates LightGBM choice.",
         ACCENT_TEAL),
        ("Runge & Zmeureanu (2019)",
         "Forecasting Energy Use in Buildings  ·  Energies 12(18), 3355",
         "Comprehensive ML survey. GBM consistently best on tabular building data; lag + rolling features most "
         "informative; CV-RMSE recommended for cross-building comparison → validates our feature + metric strategy.",
         ACCENT_GOLD),
        ("Ke et al. (2017)",
         "LightGBM: Highly Efficient GBDT  ·  NeurIPS 2017",
         "GOSS + EFB optimizations make LightGBM orders of magnitude faster than exact-split methods at 50M+ rows. "
         "Native categorical support handles primary_use and site_id without one-hot explosion.",
         ACCENT_CORAL),
        ("Molina-Solana et al. (2017)",
         "Data Science for Building Energy Mgmt  ·  Renewable & Sustainable Energy Reviews 70",
         "Residual-based anomaly detection is the industrially dominant technique (ISO 50001, ASHRAE Guideline 14). "
         "Identifies decision support as the critical missing link in academic energy ML — exactly what we build.",
         RGBColor(0xA8, 0x8E, 0xFF)),
        ("Fan et al. (2021)",
         "COVID-19 Impact on Building Energy  ·  Applied Energy 285",
         "Real-world validation: COVID occupancy collapse created massive residuals in pre-COVID baseline models, "
         "correctly flagging anomalous buildings. Validates our residual-based anomaly detection in practice.",
         RGBColor(0x4E, 0xC9, 0x71)),
        ("Chen & Guestrin (2016)",
         "XGBoost: Scalable Tree Boosting  ·  KDD 2016",
         "Baseline comparison model. LightGBM vs. XGBoost at 53M rows demonstrates concrete speedup from "
         "histogram-based splitting — a measurable, benchmarkable result in our pipeline.",
         TEXT_LIGHT),
    ]

    # 2-column layout
    for i, (auth, venue, body, col) in enumerate(papers):
        col_idx = i % 2
        row_idx = i // 2
        bx = Inches(0.5) + Inches(col_idx * 6.4)
        by = Inches(1.65) + Inches(row_idx * 1.85)
        add_rect(slide, bx, by, Inches(6.1), Inches(1.75), BG_CARD)
        add_rect(slide, bx, by, Inches(0.06), Inches(1.75), col)
        add_textbox(slide, auth, left=bx + Inches(0.15), top=by + Inches(0.1),
                    width=Inches(5.8), height=Inches(0.3),
                    font_size=13, bold=True, color=col)
        add_textbox(slide, venue, left=bx + Inches(0.15), top=by + Inches(0.38),
                    width=Inches(5.8), height=Inches(0.28),
                    font_size=10, color=TEXT_DIM, italic=True)
        add_textbox(slide, body, left=bx + Inches(0.15), top=by + Inches(0.68),
                    width=Inches(5.8), height=Inches(0.9),
                    font_size=11, color=TEXT_LIGHT)

    # Gap we address
    add_rect(slide, Inches(0.5), Inches(7.05), Inches(12.33), Inches(0.35), BG_CARD)
    add_textbox(slide,
                "GAP WE ADDRESS:  Most prior work either (a) optimizes accuracy without addressing 50M+ row scalability, or (b) scales "
                "without the full M&V workflow. Metrik AI bridges both.",
                left=Inches(0.65), top=Inches(7.08), width=Inches(12.1), height=Inches(0.3),
                font_size=11, color=ACCENT_TEAL, bold=True)

    slide_number_tag(slide, 9)
    return slide


def slide_10_roadmap(prs):
    """Project Roadmap"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, RGBColor(0xA8, 0x8E, 0xFF))

    section_label(slide, "Project Roadmap")
    slide_title(slide, "5 Phases · 14 Weeks · Measurable Milestones",
                subtitle="Front-loaded: working baseline by Week 2 · Advanced optimizations from Week 5 onward")

    phases = [
        ("Phase 1\nWks 1–2",  ACCENT_TEAL,
         "Setup + Baseline",
         ["Chunked loader · data validated",
          "Naive baseline RMSE on 5M rows",
          "Peak memory documented"]),
        ("Phase 2\nWks 3–4",  ACCENT_GOLD,
         "Features + LightGBM",
         ["Leakage-safe feature module",
          "LightGBM vs. baseline comparison",
          "cProfile first pass"]),
        ("Phase 3\nWks 5–7",  ACCENT_CORAL,
         "Anomaly + Proposal",
         ["Numba anomaly scoring built",
          "Decision-support audit list",
          "Proposal presentation (Wk 7)"]),
        ("Phase 4\nWks 8–11", RGBColor(0xA8, 0x8E, 0xFF),
         "Parallelism + GPU + Spark",
         ["Multi-worker parallel training",
          "CuPy GPU anomaly scoring",
          "PySpark prototype on 53M rows"]),
        ("Phase 5\nWks 12–14", RGBColor(0x4E, 0xC9, 0x71),
         "Polish + Report + Slides",
         ["Benchmark table complete",
          "Final report (≤4 pages)",
          "Final presentation"]),
    ]

    box_w = Inches(2.3)
    for i, (label, col, title, bullets) in enumerate(phases):
        bx = Inches(0.45) + Inches(i * 2.55)

        # Phase header
        add_rect(slide, bx, Inches(1.65), box_w, Inches(0.85), col)
        add_textbox(slide, label, left=bx, top=Inches(1.67), width=box_w, height=Inches(0.82),
                    font_size=13, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)

        # Arrow (except last)
        if i < len(phases) - 1:
            add_textbox(slide, "→",
                        left=bx + box_w, top=Inches(1.92),
                        width=Inches(0.25), height=Inches(0.35),
                        font_size=20, color=col, align=PP_ALIGN.CENTER)

        # Title band
        add_rect(slide, bx, Inches(2.5), box_w, Inches(0.35), BG_CARD)
        add_textbox(slide, title, left=bx + Inches(0.05), top=Inches(2.52),
                    width=box_w - Inches(0.1), height=Inches(0.3),
                    font_size=12, bold=True, color=col, align=PP_ALIGN.CENTER)

        # Bullet points
        add_rect(slide, bx, Inches(2.85), box_w, Inches(1.1), BG_CARD)
        for j, b in enumerate(bullets):
            add_textbox(slide, f"▸  {b}",
                        left=bx + Inches(0.08), top=Inches(2.92) + Inches(j * 0.33),
                        width=box_w - Inches(0.1), height=Inches(0.3),
                        font_size=11, color=TEXT_LIGHT)

    # Risk mitigation
    divider_line(slide, Inches(4.15))
    add_textbox(slide, "RISK MITIGATION STRATEGY",
                left=Inches(0.5), top=Inches(4.22), width=Inches(5), height=Inches(0.28),
                font_size=10, bold=True, color=ACCENT_TEAL)
    risks = [
        "Scale risk: all benchmarks on documented 5–10M subset; full 53M on NYU HPC (Greene) as stretch goal",
        "Leakage risk: unit tests assert no future timestamp in any feature — built in Phase 2, run on every commit",
        "Scope creep: anomaly detection stays residual-based only (MAD/Z-score); no unsupervised clustering",
        "Reproducibility: requirements.txt + pinned seeds + README with exact run commands from day one",
    ]
    for i, r in enumerate(risks):
        col_pos = Inches(0.5 + (i % 2) * 6.4)
        row_pos = Inches(4.55 + (i // 2) * 0.55)
        add_textbox(slide, f"▸  {r}", left=col_pos, top=row_pos,
                    width=Inches(6.1), height=Inches(0.48),
                    font_size=12, color=TEXT_LIGHT)

    slide_number_tag(slide, 10)
    return slide


def slide_11_closing(prs):
    """Why Tier-1 / Closing"""
    slide = prs.slides.add_slide(blank_layout(prs))
    add_bg(slide, BG_DARK)
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT_TEAL)
    add_rect(slide, Inches(0.08), SLIDE_H - Inches(0.04), SLIDE_W, Inches(0.04), ACCENT_GOLD)

    section_label(slide, "Summary")
    slide_title(slide, "Why Metrik AI Is a Tier-1 Project",
                subtitle="Real-world impact · Technical depth · Full course coverage · Measurable results")

    # Three pillars
    pillars = [
        ("Real-World Impact",     ACCENT_TEAL,
         ["Buildings = #1 energy efficiency opportunity in the U.S. (DOE)",
          "A production version deployed at one university → 15–25% cost savings",
          "M&V workflow is mandated by ASHRAE Guideline 14 and ISO 50001",
          "Used by Energy Services Companies (ESCOs) today"]),
        ("Technical Depth",       ACCENT_GOLD,
         ["10 of 13 course weeks directly applied as load-bearing components",
          "Benchmark table: quantitative speedup evidence, not just claims",
          "3-component output = complete system, not a single model experiment",
          "53.6M rows: the threshold where advanced Python is strictly necessary"]),
        ("Data & Rigor",          ACCENT_CORAL,
         ["Peer-reviewed dataset from real operational buildings (Miller et al. 2020)",
          "Largest open non-residential building energy benchmark in existence",
          "Leakage-safe by construction · unit-tested · reproducible",
          "Grounded in 6 published papers across top venues"]),
    ]

    for i, (title, col, pts) in enumerate(pillars):
        bx = Inches(0.5) + Inches(i * 4.15)
        add_rect(slide, bx, Inches(1.65), Inches(3.95), Inches(0.38), col)
        add_textbox(slide, title, left=bx, top=Inches(1.67), width=Inches(3.95), height=Inches(0.35),
                    font_size=14, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        add_rect(slide, bx, Inches(2.03), Inches(3.95), Inches(1.75), BG_CARD)
        for j, pt in enumerate(pts):
            add_textbox(slide, f"▸  {pt}",
                        left=bx + Inches(0.12), top=Inches(2.1) + Inches(j * 0.4),
                        width=Inches(3.7), height=Inches(0.37),
                        font_size=12, color=TEXT_LIGHT)

    # Bloomberg analogy box
    add_rect(slide, Inches(0.5), Inches(4.0), Inches(12.33), Inches(0.85), BG_CARD)
    add_rect(slide, Inches(0.5), Inches(4.0), Inches(0.06), Inches(0.85), ACCENT_GOLD)
    add_textbox(slide,
                '"Metrik AI does for building energy what a Bloomberg terminal does for financial markets — '
                'it turns raw, messy, 53-million-row data into actionable operational intelligence, '
                'at scale, in real time."',
                left=Inches(0.7), top=Inches(4.08), width=Inches(11.9), height=Inches(0.72),
                font_size=15, italic=True, color=TEXT_WHITE, align=PP_ALIGN.LEFT)

    # 5-min talk timer strip
    divider_line(slide, Inches(5.05))
    add_textbox(slide, "5-MINUTE TALK FLOW",
                left=Inches(0.5), top=Inches(5.12), width=Inches(3), height=Inches(0.28),
                font_size=10, bold=True, color=ACCENT_TEAL)
    talk_flow = [
        ("0:00–0:45", "Motivation — 40% energy, $130B waste, stakeholders"),
        ("0:45–1:15", "Problem — formal statement, 3-part decomposition"),
        ("1:15–2:00", "Solution — pipeline architecture, 3 components"),
        ("2:00–2:45", "Data — ASHRAE GEPIII, schema, features"),
        ("2:45–3:45", "Course techniques — 10 weeks mapped to pipeline"),
        ("3:45–4:15", "Literature — 6 papers, gap we address"),
        ("4:15–5:00", "Roadmap + closing argument"),
    ]
    for i, (t, desc) in enumerate(talk_flow):
        bx_off = Inches(0.5 + (i % 4) * 3.1) if i < 4 else Inches(0.5 + (i - 4) * 3.85)
        by_off = Inches(5.46) if i < 4 else Inches(6.1)
        add_rect(slide, bx_off, by_off, Inches(2.95), Inches(0.5), BG_CARD)
        add_textbox(slide, t, left=bx_off + Inches(0.05), top=by_off + Inches(0.04),
                    width=Inches(0.9), height=Inches(0.4),
                    font_size=10, bold=True, color=ACCENT_TEAL)
        add_textbox(slide, desc, left=bx_off + Inches(1.0), top=by_off + Inches(0.06),
                    width=Inches(1.85), height=Inches(0.38),
                    font_size=10, color=TEXT_LIGHT)

    slide_number_tag(slide, 11)
    return slide


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    print("Building slides...")
    slide_01_title(prs);            print("  [1/11] Title")
    slide_02_motivation(prs);       print("  [2/11] Real-World Motivation")
    slide_03_problem(prs);          print("  [3/11] Problem Statement")
    slide_04_solution(prs);         print("  [4/11] Proposed Solution")
    slide_05_dataset_overview(prs); print("  [5/11] Dataset Overview")
    slide_06_dataset_deep(prs);     print("  [6/11] Dataset Deep Dive")
    slide_07_technical_approach(prs); print("  [7/11] Technical Approach")
    slide_08_course_techniques(prs); print("  [8/11] Course Techniques")
    slide_09_literature(prs);       print("  [9/11] Literature Review")
    slide_10_roadmap(prs);          print(" [10/11] Roadmap")
    slide_11_closing(prs);          print(" [11/11] Closing / Why Tier-1")

    out = "/Users/akhileshvangala/Desktop/Metrik AI/Proposal Files/MetrikAI_Proposal_Presentation.pptx"
    prs.save(out)
    print(f"\nSaved → {out}")


if __name__ == "__main__":
    main()
